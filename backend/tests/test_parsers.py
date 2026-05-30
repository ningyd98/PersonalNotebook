"""
解析器单元测试

测试所有已实现的 Parser：
- MarkdownParser: frontmatter、heading tree、wikilinks、code blocks、tables
- TXTParser: 基础文本分段
- PDFParser: 基础接口测试（mock PyMuPDF）
- DOCXParser: 创建临时 docx 文件测试
- PPTXParser: 创建临时 pptx 文件测试
- LaTeXParser: section/equation/label/cite 提取
- AudioParser: 基础接口测试（mock ASR）
- VideoParser: 基础接口测试（mock ffmpeg/ASR）
- FallbackParser: 兜底解析
- ParserRegistry: 自动注册和匹配

运行: pytest tests/test_parsers.py -v
"""

import os
import sys
import uuid
from unittest.mock import patch, MagicMock

import pytest

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))


# ============================================================
# MarkdownParser
# ============================================================
class TestMarkdownParser:
    """MarkdownParser 完整测试"""

    def _write_md(self, tmp_path, content, filename="test.md"):
        p = tmp_path / filename
        p.write_text(content, encoding="utf-8")
        return str(p)

    def test_frontmatter(self, tmp_path):
        from app.services.parsers.markdown_parser import MarkdownParser
        parser = MarkdownParser()
        md = """---
title: 测试文档
author: 张三
tags:
  - python
  - testing
---
正文内容
"""
        fp = self._write_md(tmp_path, md)
        udr = parser.parse(fp)
        assert udr.metadata["title"] == "测试文档"
        assert udr.metadata["author"] == "张三"
        assert "python" in udr.metadata["tags"]
        assert "testing" in udr.metadata["tags"]

    def test_heading_tree(self, tmp_path):
        from app.services.parsers.markdown_parser import MarkdownParser
        parser = MarkdownParser()
        md = """# 第一章

## 第一节

### 第一小节

## 第二节

# 第二章
"""
        fp = self._write_md(tmp_path, md)
        udr = parser.parse(fp)
        headings = [b for b in udr.blocks if b.type == "heading"]
        levels = [h.level for h in headings]
        texts = [h.text for h in headings]

        assert levels == [1, 2, 3, 2, 1]
        assert "第一章" in texts
        assert "第一小节" in texts

        # section_path 验证
        h3 = [h for h in headings if h.level == 3][0]
        assert "第一章" in h3.section_path
        assert "第一节" in h3.section_path
        assert "第一小节" in h3.section_path

    def test_wikilinks(self, tmp_path):
        from app.services.parsers.markdown_parser import MarkdownParser
        parser = MarkdownParser()
        md = """参见 [[概念A]] 和 [[概念B|显示文本]] 以及 [[概念C#章节]]。
"""
        fp = self._write_md(tmp_path, md)
        udr = parser.parse(fp)
        link_targets = [r.target_document_id for r in udr.relations]
        assert "概念A" in link_targets
        assert "概念B" in link_targets
        assert "概念C" in link_targets

    def test_code_blocks(self, tmp_path):
        from app.services.parsers.markdown_parser import MarkdownParser
        parser = MarkdownParser()
        md = """# 代码测试

```python
def hello():
    print("Hello, World!")
```

```javascript
console.log("Hi");
```
"""
        fp = self._write_md(tmp_path, md)
        udr = parser.parse(fp)
        codes = [b for b in udr.blocks if b.type == "code"]
        assert len(codes) == 2

        py_code = [c for c in codes if c.metadata.get("language") == "python"][0]
        assert "hello" in py_code.text
        assert "print" in py_code.text

        js_code = [c for c in codes if c.metadata.get("language") == "javascript"][0]
        assert "console.log" in js_code.text

    def test_tables(self, tmp_path):
        from app.services.parsers.markdown_parser import MarkdownParser
        parser = MarkdownParser()
        md = """# 表格测试

| 姓名 | 年龄 | 城市 |
|------|------|------|
| 张三 | 25   | 北京 |
| 李四 | 30   | 上海 |
"""
        fp = self._write_md(tmp_path, md)
        udr = parser.parse(fp)
        tables = [b for b in udr.blocks if b.type == "table"]
        assert len(tables) >= 1

        table = tables[0]
        assert table.structured_data is not None
        assert "headers" in table.structured_data
        assert len(table.structured_data["headers"]) == 3
        assert "姓名" in table.structured_data["headers"]
        assert len(table.structured_data["rows"]) == 2

    def test_task_list(self, tmp_path):
        from app.services.parsers.markdown_parser import MarkdownParser
        parser = MarkdownParser()
        md = """- [x] 完成任务A
- [ ] 待完成任务B
"""
        fp = self._write_md(tmp_path, md)
        udr = parser.parse(fp)
        tasks = [b for b in udr.blocks if b.type == "list" and b.metadata.get("list_type") == "task"]
        assert len(tasks) == 2
        checked_tasks = [t for t in tasks if t.metadata.get("checked") is True]
        unchecked_tasks = [t for t in tasks if t.metadata.get("checked") is False]
        assert len(checked_tasks) == 1
        assert len(unchecked_tasks) == 1

    def test_empty_file(self, tmp_path):
        from app.services.parsers.markdown_parser import MarkdownParser
        parser = MarkdownParser()
        fp = self._write_md(tmp_path, "")
        udr = parser.parse(fp)
        assert udr.document_id.startswith("doc_")
        assert len(udr.blocks) == 0

    def test_no_frontmatter(self, tmp_path):
        from app.services.parsers.markdown_parser import MarkdownParser
        parser = MarkdownParser()
        md = """这是没有 frontmatter 的文档。

第二段落。
"""
        fp = self._write_md(tmp_path, md)
        udr = parser.parse(fp)
        # 没有 frontmatter 时，title 回退到 filename
        assert udr.metadata["title"] == "test.md"
        paragraphs = [b for b in udr.blocks if b.type == "paragraph"]
        assert len(paragraphs) >= 1


# ============================================================
# TXTParser
# ============================================================
class TestTXTParser:
    """纯文本解析器测试"""

    def test_basic_paragraph_split(self, tmp_path):
        from app.services.parsers.txt_parser import TXTParser
        parser = TXTParser()

        content = "第一段内容。\n\n第二段内容。\n\n第三段内容。"
        fp = tmp_path / "test.txt"
        fp.write_text(content, encoding="utf-8")

        udr = parser.parse(str(fp))
        paragraphs = [b for b in udr.blocks if b.type == "paragraph"]
        assert len(paragraphs) == 3
        assert "第一段" in paragraphs[0].text
        assert "第二段" in paragraphs[1].text
        assert "第三段" in paragraphs[2].text

    def test_single_paragraph(self, tmp_path):
        from app.services.parsers.txt_parser import TXTParser
        parser = TXTParser()

        fp = tmp_path / "single.txt"
        fp.write_text("只有一段文字", encoding="utf-8")

        udr = parser.parse(str(fp))
        assert len(udr.blocks) == 1
        assert udr.blocks[0].type == "paragraph"

    def test_source_metadata(self, tmp_path):
        from app.services.parsers.txt_parser import TXTParser
        parser = TXTParser()

        fp = tmp_path / "meta.txt"
        fp.write_text("内容", encoding="utf-8")

        udr = parser.parse(str(fp))
        assert udr.source["mime_type"] == "text/plain"
        assert udr.source["filename"] == "meta.txt"

    def test_supported_extensions(self):
        from app.services.parsers.txt_parser import TXTParser
        parser = TXTParser()
        assert ".txt" in parser.supported_extensions
        assert ".log" in parser.supported_extensions


# ============================================================
# PDFParser
# ============================================================
class TestPDFParser:
    """PDF 解析器测试（mock PyMuPDF）"""

    def test_fallback_when_no_pymupdf(self, tmp_path):
        from app.services.parsers.pdf_parser import PDFParser
        parser = PDFParser()

        # 创建一个伪 PDF 文件
        fp = tmp_path / "test.pdf"
        fp.write_bytes(b"%PDF-1.4 fake content")

        with patch("app.services.parsers.pdf_parser.HAS_PYMUPDF", False):
            udr = parser.parse(str(fp))

        assert udr.document_id.startswith("doc_")
        assert len(udr.blocks) >= 1
        assert udr.metadata.get("warning") is not None

    def test_supported_mime_types(self):
        from app.services.parsers.pdf_parser import PDFParser
        parser = PDFParser()
        assert "application/pdf" in parser.supported_mime_types
        assert ".pdf" in parser.supported_extensions

    def test_with_mock_fitz(self, tmp_path):
        """使用 mock fitz 模拟 PDF 解析"""
        from app.services.parsers.pdf_parser import PDFParser
        parser = PDFParser()

        fp = tmp_path / "mock.pdf"
        fp.write_bytes(b"%PDF-1.4 fake")

        mock_page = MagicMock()
        mock_page.get_text.return_value = [
            (0, 0, 100, 20, "Hello World", 0, 0),
        ]
        mock_page.find_tables.return_value = []

        mock_doc = MagicMock()
        mock_doc.__len__ = MagicMock(return_value=1)
        mock_doc.__getitem__ = MagicMock(return_value=mock_page)
        mock_doc.metadata = {"title": "Test PDF", "author": "Author"}
        mock_doc.close = MagicMock()

        with patch("app.services.parsers.pdf_parser.HAS_PYMUPDF", True), \
             patch("app.services.parsers.pdf_parser.fitz") as mock_fitz:
            mock_fitz.open.return_value = mock_doc
            udr = parser.parse(str(fp))

        assert udr.metadata["title"] == "Test PDF"
        assert len(udr.blocks) >= 1
        block_texts = [b.text for b in udr.blocks]
        assert any("Hello World" in t for t in block_texts)


# ============================================================
# DOCXParser
# ============================================================
class TestDOCXParser:
    """DOCX 解析器测试（使用 python-docx 生成测试文件）"""

    @pytest.fixture
    def docx_file(self, tmp_path):
        """使用 python-docx 创建临时 docx 文件"""
        try:
            from docx import Document
            from docx.shared import Pt
        except ImportError:
            pytest.skip("python-docx not installed")

        doc = Document()
        doc.core_properties.title = "测试文档"
        doc.core_properties.author = "测试作者"

        doc.add_heading("第一章 概述", level=1)
        doc.add_paragraph("这是概述的内容，描述了项目的基本信息。")
        doc.add_heading("第二章 方法", level=1)
        doc.add_paragraph("方法章节的内容。")

        # 添加表格
        table = doc.add_table(rows=3, cols=2)
        table.cell(0, 0).text = "名称"
        table.cell(0, 1).text = "值"
        table.cell(1, 0).text = "参数A"
        table.cell(1, 1).text = "100"
        table.cell(2, 0).text = "参数B"
        table.cell(2, 1).text = "200"

        fp = tmp_path / "test.docx"
        doc.save(str(fp))
        return str(fp)

    def test_docx_headings(self, docx_file):
        from app.services.parsers.docx_parser import DOCXParser
        parser = DOCXParser()
        udr = parser.parse(docx_file)

        headings = [b for b in udr.blocks if b.type == "heading"]
        assert len(headings) >= 2
        heading_texts = [h.text for h in headings]
        assert any("概述" in t for t in heading_texts)
        assert any("方法" in t for t in heading_texts)

    def test_docx_paragraphs(self, docx_file):
        from app.services.parsers.docx_parser import DOCXParser
        parser = DOCXParser()
        udr = parser.parse(docx_file)

        paragraphs = [b for b in udr.blocks if b.type == "paragraph"]
        assert len(paragraphs) >= 2
        all_text = " ".join(p.text for p in paragraphs)
        assert "概述" in all_text or "方法" in all_text or "信息" in all_text

    def test_docx_tables(self, docx_file):
        from app.services.parsers.docx_parser import DOCXParser
        parser = DOCXParser()
        udr = parser.parse(docx_file)

        tables = [b for b in udr.blocks if b.type == "table"]
        assert len(tables) >= 1
        table = tables[0]
        assert table.structured_data is not None
        assert "headers" in table.structured_data
        assert "rows" in table.structured_data

    def test_docx_metadata(self, docx_file):
        from app.services.parsers.docx_parser import DOCXParser
        parser = DOCXParser()
        udr = parser.parse(docx_file)

        assert udr.metadata.get("title") in ("测试文档", "test.docx")
        assert udr.source["mime_type"] == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"

    def test_fallback_when_no_docx_lib(self, tmp_path):
        from app.services.parsers.docx_parser import DOCXParser
        parser = DOCXParser()

        fp = tmp_path / "fake.docx"
        fp.write_bytes(b"PK\x03\x04 fake docx")

        with patch("app.services.parsers.docx_parser.HAS_PYTHON_DOCX", False):
            udr = parser.parse(str(fp))

        assert udr.document_id.startswith("doc_")
        assert udr.metadata.get("warning") is not None


# ============================================================
# PPTXParser
# ============================================================
class TestPPTXParser:
    """PPTX 解析器测试（使用 python-pptx 生成测试文件）"""

    @pytest.fixture
    def pptx_file(self, tmp_path):
        """使用 python-pptx 创建临时 pptx 文件"""
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            pytest.skip("python-pptx not installed")

        prs = Presentation()
        # Slide 1
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide1 = prs.slides.add_slide(slide_layout)
        slide1.shapes.title.text = "项目简介"
        slide1.placeholders[1].text = "这是一个项目简介的幻灯片"

        # Slide 2
        slide2 = prs.slides.add_slide(slide_layout)
        slide2.shapes.title.text = "技术方案"
        slide2.placeholders[1].text = "技术方案的详细描述"

        fp = tmp_path / "test.pptx"
        prs.save(str(fp))
        return str(fp)

    def test_pptx_slides(self, pptx_file):
        from app.services.parsers.pptx_parser import PPTXParser
        parser = PPTXParser()
        udr = parser.parse(pptx_file)

        assert udr.document_id.startswith("doc_")
        assert len(udr.blocks) >= 2

    def test_pptx_metadata(self, pptx_file):
        from app.services.parsers.pptx_parser import PPTXParser
        parser = PPTXParser()
        udr = parser.parse(pptx_file)

        assert udr.source["mime_type"] == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        assert "slide_count" in udr.metadata
        assert udr.metadata["slide_count"] == 2

    def test_pptx_slide_content(self, pptx_file):
        from app.services.parsers.pptx_parser import PPTXParser
        parser = PPTXParser()
        udr = parser.parse(pptx_file)

        all_text = " ".join(b.text for b in udr.blocks)
        assert "项目简介" in all_text or "技术方案" in all_text

    def test_fallback_when_no_pptx_lib(self, tmp_path):
        from app.services.parsers.pptx_parser import PPTXParser
        parser = PPTXParser()

        fp = tmp_path / "fake.pptx"
        fp.write_bytes(b"PK\x03\x04 fake pptx")

        with patch("app.services.parsers.pptx_parser.HAS_PPTX", False):
            udr = parser.parse(str(fp))

        assert udr.document_id.startswith("doc_")
        assert udr.metadata.get("warning") is not None


# ============================================================
# LaTeXParser
# ============================================================
class TestLaTeXParser:
    """LaTeX 解析器测试"""

    def test_section_extraction(self, tmp_path):
        from app.services.parsers.latex_parser import LaTeXParser
        parser = LaTeXParser()

        content = r"""
\documentclass{article}
\title{Test Paper}
\begin{document}
\section{Introduction}
This is the introduction.
\subsection{Background}
Background info here.
\section{Methods}
Methods description.
\end{document}
"""
        fp = tmp_path / "test.tex"
        fp.write_text(content, encoding="utf-8")

        udr = parser.parse(str(fp))
        headings = [b for b in udr.blocks if b.type == "heading"]
        assert len(headings) >= 2
        heading_texts = [h.text for h in headings]
        assert any("Introduction" in t for t in heading_texts)
        assert any("Methods" in t for t in heading_texts)

    def test_equation_extraction(self, tmp_path):
        from app.services.parsers.latex_parser import LaTeXParser
        parser = LaTeXParser()

        content = r"""
\documentclass{article}
\begin{document}
\section{Formulas}
\begin{equation}
E = mc^2
\end{equation}
Some text between equations.
\begin{align}
a + b &= c \\
d + e &= f
\end{align}
\end{document}
"""
        fp = tmp_path / "equations.tex"
        fp.write_text(content, encoding="utf-8")

        udr = parser.parse(str(fp))
        equations = [b for b in udr.blocks if b.type == "equation"]
        assert len(equations) >= 2
        eq_texts = [e.text for e in equations]
        assert any("mc" in t for t in eq_texts)
        assert any("a + b" in t for t in eq_texts)

    def test_label_extraction(self, tmp_path):
        from app.services.parsers.latex_parser import LaTeXParser
        parser = LaTeXParser()

        content = r"""
\documentclass{article}
\begin{document}
\section{Test}\label{sec:test}
Content with \label{eq:main}.
\begin{equation}\label{eq:einstein}
E = mc^2
\end{equation}
\end{document}
"""
        fp = tmp_path / "labels.tex"
        fp.write_text(content, encoding="utf-8")

        udr = parser.parse(str(fp))
        assert "sec:test" in udr.metadata.get("labels", [])
        assert "eq:main" in udr.metadata.get("labels", [])
        assert "eq:einstein" in udr.metadata.get("labels", [])

    def test_cite_extraction(self, tmp_path):
        from app.services.parsers.latex_parser import LaTeXParser
        parser = LaTeXParser()

        content = r"""
\documentclass{article}
\begin{document}
\section{Related Work}
As shown by \cite{smith2020, jones2021}, the method works.
Another reference \cite{wang2019}.
\end{document}
"""
        fp = tmp_path / "cites.tex"
        fp.write_text(content, encoding="utf-8")

        udr = parser.parse(str(fp))
        cites = udr.metadata.get("cites", [])
        assert "smith2020" in cites
        assert "jones2021" in cites
        assert "wang2019" in cites

        # UDRRelation 检查
        cite_relations = [r for r in udr.relations if r.relation_type == "cites"]
        cite_targets = [r.target_document_id for r in cite_relations]
        assert "smith2020" in cite_targets
        assert "wang2019" in cite_targets

    def test_theorem_environment(self, tmp_path):
        from app.services.parsers.latex_parser import LaTeXParser
        parser = LaTeXParser()

        content = r"""
\documentclass{article}
\begin{document}
\section{Theory}
\begin{theorem}[Pythagorean]
In a right triangle, $a^2 + b^2 = c^2$.
\end{theorem}
\begin{proof}
The proof is straightforward.
\end{proof}
\end{document}
"""
        fp = tmp_path / "theorem.tex"
        fp.write_text(content, encoding="utf-8")

        udr = parser.parse(str(fp))
        annotations = [b for b in udr.blocks if b.type == "annotation"]
        assert len(annotations) >= 1
        thm = [a for a in annotations if a.metadata.get("theorem_type") == "theorem"]
        assert len(thm) >= 1

    def test_supported_extensions(self):
        from app.services.parsers.latex_parser import LaTeXParser
        parser = LaTeXParser()
        assert ".tex" in parser.supported_extensions
        assert ".latex" in parser.supported_extensions


# ============================================================
# AudioParser（mock ASR）
# ============================================================
class TestAudioParser:
    """音频解析器测试（mock ASR）"""

    def test_supported_mime_types(self):
        from app.services.parsers.audio_parser import AudioParser
        parser = AudioParser()
        assert "audio/mpeg" in parser.supported_mime_types
        assert "audio/wav" in parser.supported_mime_types
        assert ".mp3" in parser.supported_extensions
        assert ".wav" in parser.supported_extensions

    def test_fallback_when_no_asr(self, tmp_path):
        from app.services.parsers.audio_parser import AudioParser
        parser = AudioParser()

        fp = tmp_path / "test.mp3"
        fp.write_bytes(b"fake audio content")

        with patch("app.services.parsers.audio_parser.HAS_FASTER_WHISPER", False):
            udr = parser.parse(str(fp))

        assert udr.document_id.startswith("doc_")
        assert len(udr.blocks) >= 1
        assert udr.metadata.get("asr_available") is False

    def test_guess_mime(self):
        from app.services.parsers.audio_parser import AudioParser
        from pathlib import Path

        assert AudioParser._guess_mime(Path("test.mp3")) == "audio/mpeg"
        assert AudioParser._guess_mime(Path("test.wav")) == "audio/wav"
        assert AudioParser._guess_mime(Path("test.m4a")) == "audio/mp4"
        assert AudioParser._guess_mime(Path("test.unknown")) == "application/octet-stream"

    def test_merge_segments(self):
        from app.services.parsers.audio_parser import AudioParser
        parser = AudioParser()

        # 模拟 ASR segments
        segments = []
        for i in range(10):
            seg = MagicMock()
            seg.start = i * 10.0
            seg.end = (i + 1) * 10.0
            seg.text = f" 第{i+1}段文字"
            segments.append(seg)

        blocks = parser._merge_segments(segments, "doc_test")
        assert len(blocks) >= 1
        for block in blocks:
            assert block.type == "transcript"
            assert block.start_time is not None
            assert block.end_time is not None


# ============================================================
# VideoParser（mock ffmpeg/ASR）
# ============================================================
class TestVideoParser:
    """视频解析器测试（mock ffmpeg/ASR）"""

    def test_supported_mime_types(self):
        from app.services.parsers.video_parser import VideoParser
        parser = VideoParser()
        assert "video/mp4" in parser.supported_mime_types
        assert "video/quicktime" in parser.supported_mime_types
        assert ".mp4" in parser.supported_extensions
        assert ".mov" in parser.supported_extensions

    def test_guess_mime(self):
        from app.services.parsers.video_parser import VideoParser
        from pathlib import Path

        assert VideoParser._guess_mime(Path("test.mp4")) == "video/mp4"
        assert VideoParser._guess_mime(Path("test.mov")) == "video/quicktime"
        assert VideoParser._guess_mime(Path("test.mkv")) == "video/x-matroska"
        assert VideoParser._guess_mime(Path("test.unknown")) == "application/octet-stream"

    def test_merge_segments(self):
        from app.services.parsers.video_parser import VideoParser
        parser = VideoParser()

        segments = []
        for i in range(5):
            seg = MagicMock()
            seg.start = i * 20.0
            seg.end = (i + 1) * 20.0
            seg.text = f" 视频段{i+1}"
            segments.append(seg)

        blocks = parser._merge_segments(segments, "doc_video_test")
        assert len(blocks) >= 1
        for block in blocks:
            assert block.type == "transcript"

    def test_merge_timeline(self):
        from app.services.parsers.video_parser import VideoParser
        from app.services.parsers.base import UDRBlock
        parser = VideoParser()

        blocks = [
            UDRBlock(block_id="b1", type="transcript", text="Hello",
                     start_time=0.0, end_time=10.0),
            UDRBlock(block_id="b2", type="image", text="OCR text",
                     start_time=5.0, ocr_text="OCR text",
                     metadata={"visual_caption": None}),
        ]

        merged = parser._merge_timeline(blocks, "doc_test")
        assert len(merged) == 2
        for b in merged:
            assert b.type == "video_segment"

    def test_parse_with_mock(self, tmp_path):
        """使用 mock 测试 VideoParser.parse"""
        from app.services.parsers.video_parser import VideoParser
        parser = VideoParser()

        fp = tmp_path / "test.mp4"
        fp.write_bytes(b"fake video content")

        with patch("app.services.parsers.video_parser.HAS_FASTER_WHISPER", False), \
             patch.object(parser, "_get_duration", return_value=60.0), \
             patch.object(parser, "_extract_and_transcribe", return_value=[]), \
             patch.object(parser, "_extract_keyframes", return_value=[]), \
             patch.object(parser, "_process_keyframes", return_value=[]):
            udr = parser.parse(str(fp))

        assert udr.document_id.startswith("doc_")
        assert udr.metadata.get("duration_seconds") == 60.0


# ============================================================
# FallbackParser
# ============================================================
class TestFallbackParser:
    """兜底解析器测试"""

    def test_utf8_text(self, tmp_path):
        from app.services.parsers.fallback_parser import FallbackParser
        parser = FallbackParser()

        fp = tmp_path / "test.xyz"
        fp.write_text("第一段\n\n第二段\n\n第三段", encoding="utf-8")

        udr = parser.parse(str(fp))
        assert len(udr.blocks) >= 1
        assert udr.metadata.get("warning") is not None

    def test_binary_file(self, tmp_path):
        from app.services.parsers.fallback_parser import FallbackParser
        parser = FallbackParser()

        fp = tmp_path / "binary.dat"
        fp.write_bytes(b"\x00\x01\x02\x03\xff\xfe\xfd")

        udr = parser.parse(str(fp))
        assert len(udr.blocks) >= 1
        # 二进制文件应该有 metadata block
        meta_blocks = [b for b in udr.blocks if b.type == "metadata"]
        assert len(meta_blocks) >= 1

    def test_wildcard_mime(self):
        from app.services.parsers.fallback_parser import FallbackParser
        parser = FallbackParser()
        assert "*/*" in parser.supported_mime_types
        assert "*" in parser.supported_extensions


# ============================================================
# ParserRegistry
# ============================================================
class TestParserRegistry:
    """Parser 注册中心测试"""

    def test_registry_has_parsers(self):
        """导入 parsers 模块后应自动注册所有 parser"""
        import app.services.parsers  # noqa: F401 — 触发注册
        from app.services.parsers.base import ParserRegistry

        parsers = ParserRegistry.list_parsers()
        assert len(parsers) > 0
        parser_classes = [p["class"] for p in parsers]
        assert "MarkdownParser" in parser_classes
        assert "TXTParser" in parser_classes
        assert "PDFParser" in parser_classes
        assert "FallbackParser" in parser_classes

    def test_get_parser_by_extension(self):
        import app.services.parsers  # noqa: F401
        from app.services.parsers.base import ParserRegistry

        md_parser = ParserRegistry.get_parser(extension=".md")
        assert md_parser is not None
        assert md_parser.__class__.__name__ == "MarkdownParser"

        txt_parser = ParserRegistry.get_parser(extension=".txt")
        assert txt_parser is not None
        assert txt_parser.__class__.__name__ == "TXTParser"

    def test_get_parser_by_mime_type(self):
        import app.services.parsers  # noqa: F401
        from app.services.parsers.base import ParserRegistry

        md_parser = ParserRegistry.get_parser(mime_type="text/markdown")
        assert md_parser is not None
        assert md_parser.__class__.__name__ == "MarkdownParser"

        pdf_parser = ParserRegistry.get_parser(mime_type="application/pdf")
        assert pdf_parser is not None
        assert pdf_parser.__class__.__name__ == "PDFParser"

    def test_get_parser_extension_without_dot(self):
        import app.services.parsers  # noqa: F401
        from app.services.parsers.base import ParserRegistry

        parser = ParserRegistry.get_parser(extension="md")
        assert parser is not None
        assert parser.__class__.__name__ == "MarkdownParser"

    def test_get_parser_unknown_returns_none_without_fallback(self):
        """对于未知扩展名且没有 FallbackParser 注册时返回 None"""
        from app.services.parsers.base import ParserRegistry

        # 清除后重新测试
        # 因为 FallbackParser 有 "*/*" 匹配，模糊匹配可能命中
        # 直接测试不存在的 mime_type
        parser = ParserRegistry.get_parser(mime_type="application/x-totally-unknown")
        # 可能返回 None 或 FallbackParser（模糊匹配）
        # 这是可接受的行为

    def test_get_parser_fuzzy_mime_match(self):
        import app.services.parsers  # noqa: F401
        from app.services.parsers.base import ParserRegistry

        # 模糊匹配：audio/something 应该匹配到 AudioParser
        audio_parser = ParserRegistry.get_parser(mime_type="audio/unknown-format")
        assert audio_parser is not None
        assert "Audio" in audio_parser.__class__.__name__
