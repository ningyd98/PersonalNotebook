"""PPTX Parser — 使用 python-pptx 解析 PowerPoint 文档"""

import uuid
from pathlib import Path
from typing import Optional

from loguru import logger

from app.services.parsers.base import (
    BaseParser, ParserRegistry, UDRAsset, UDRBlock, UnifiedDocument,
)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    logger.warning("python-pptx not installed; PPTX parsing will be limited")


class PptxParser(BaseParser):
    supported_mime_types = [
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ]
    supported_extensions = [".pptx"]

    def parse(self, file_path: str, options: Optional[dict] = None) -> UnifiedDocument:
        if not HAS_PPTX:
            return self._fallback_parse(file_path)

        doc_id = f"doc_{uuid.uuid4().hex[:12]}"
        path = Path(file_path)
        filename = path.name
        options = options or {}

        blocks: list[UDRBlock] = []
        assets: list[UDRAsset] = []

        try:
            prs = Presentation(file_path)

            for slide_idx, slide in enumerate(prs.slides):
                # Slide heading
                slide_title = f"Slide {slide_idx + 1}"
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.shape_type == 14:  # MSO_SHAPE_TYPE.TITLE
                        slide_title = shape.text_frame.text.strip() or slide_title
                        break

                blocks.append(UDRBlock(
                    block_id=f"{doc_id}_b{len(blocks):04d}",
                    type="heading",
                    text=slide_title,
                    level=1,
                    slide=slide_idx + 1,
                ))

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if not text:
                                continue
                            # Determine if heading
                            level = None
                            if para.level > 0:
                                level = para.level + 1
                            blocks.append(UDRBlock(
                                block_id=f"{doc_id}_b{len(blocks):04d}",
                                type="heading" if level else "paragraph",
                                text=text,
                                level=level,
                                slide=slide_idx + 1,
                            ))

                    if shape.shape_type == 13:  # Picture
                        try:
                            image = shape.image
                            asset_id = f"asset_{uuid.uuid4().hex[:8]}"
                            content_type = image.content_type
                            ext = content_type.split("/")[-1] if "/" in content_type else "png"
                            if ext == "jpeg":
                                ext = "jpg"
                            object_name = f"parsed_assets/{doc_id}/{asset_id}.{ext}"

                            try:
                                from app.services.storage import minio_storage
                                minio_storage.upload_file(
                                    object_name=object_name,
                                    data=image.blob,
                                    content_type=content_type,
                                )
                                asset_uri = f"minio://{minio_storage.client._base_url.host}/{object_name}"
                            except Exception:
                                asset_uri = f"inline://{asset_id}"

                            import hashlib
                            checksum = hashlib.md5(image.blob).hexdigest()
                            asset = UDRAsset(
                                asset_id=asset_id,
                                asset_type="image",
                                asset_uri=asset_uri,
                                mime_type=content_type,
                                file_size=len(image.blob),
                                checksum=checksum,
                                metadata={"source": "pptx_embedded", "slide": slide_idx + 1},
                            )
                            assets.append(asset)
                            blocks.append(UDRBlock(
                                block_id=f"{doc_id}_b{len(blocks):04d}",
                                type="image",
                                text=f"[Image: {asset_id}]",
                                slide=slide_idx + 1,
                                asset_uri=asset_uri,
                                metadata={"asset_id": asset_id},
                            ))
                        except Exception as e:
                            logger.debug(f"Failed to extract image from slide {slide_idx}: {e}")

                    if shape.has_table:
                        try:
                            table = shape.table
                            rows_data = []
                            for row in table.rows:
                                cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                                rows_data.append(cells)

                            if rows_data:
                                headers = rows_data[0]
                                data_rows = rows_data[1:] if len(rows_data) > 1 else []
                                md_lines = ["| " + " | ".join(headers) + " |"]
                                md_lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
                                for row in data_rows:
                                    padded = row + [""] * (len(headers) - len(row))
                                    md_lines.append("| " + " | ".join(padded[:len(headers)]) + " |")

                                blocks.append(UDRBlock(
                                    block_id=f"{doc_id}_b{len(blocks):04d}",
                                    type="table",
                                    text="\n".join(md_lines),
                                    slide=slide_idx + 1,
                                    structured_data={"headers": headers, "rows": data_rows},
                                ))
                        except Exception as e:
                            logger.debug(f"Failed to extract table from slide {slide_idx}: {e}")

                    # Speaker notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        blocks.append(UDRBlock(
                            block_id=f"{doc_id}_b{len(blocks):04d}",
                            type="annotation",
                            text=notes_text,
                            slide=slide_idx + 1,
                            metadata={"note_type": "speaker_notes"},
                        ))

            return UnifiedDocument(
                document_id=doc_id,
                source={
                    "filename": filename,
                    "mime_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    "source_uri": str(path.absolute()),
                },
                metadata={"title": filename, "slide_count": len(prs.slides)},
                blocks=blocks,
                assets=assets,
            )

        except Exception as e:
            logger.error(f"PPTX parsing failed: {e}")
            return self._fallback_parse(file_path)

    @staticmethod
    def _fallback_parse(file_path: str) -> UnifiedDocument:
        doc_id = f"doc_{uuid.uuid4().hex[:12]}"
        path = Path(file_path)
        return UnifiedDocument(
            document_id=doc_id,
            source={
                "filename": path.name,
                "mime_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "source_uri": str(path.absolute()),
            },
            metadata={"title": path.name, "warning": "PPTX parsing library not available"},
            blocks=[
                UDRBlock(
                    block_id=f"{doc_id}_b0000",
                    type="paragraph",
                    text=f"[PPTX file: {path.name} — install python-pptx for full parsing]",
                ),
            ],
        )


# 注册
ParserRegistry.register(PptxParser())
